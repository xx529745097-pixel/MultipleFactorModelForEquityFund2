import copy
import html
from io import BytesIO, IOBase
from typing import BinaryIO, List, Literal, Union

import altair as alt
import dataframe_image as dfi
import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pandas.io.formats.style import Styler
from pptx.presentation import Presentation as PrsCls
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
import re


def replace_text(ppt: Union[PrsCls, Slide], search_pattern: str, repl: str) -> None:
    """search and replace text in PowerPoint while preserving formatting

    Args:
        ppt: 输入的ppt文件, Presentation对象 / slide对象
        search_pattern: 查找的字符串
        repl: 替换的字符串

    Returns:
        ppt对象
    """
    # Useful Links ;)
    # https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    # https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    if isinstance(ppt, PrsCls):
        for slide in ppt.slides:
            _replace_text_in_slide(slide, search_pattern, repl)
    elif isinstance(ppt, Slide):
        _replace_text_in_slide(ppt, search_pattern, repl)


def _replace_text_in_slide(slide: Slide, search_str: str, repl: str) -> Slide:
    """Replace text within a page of ppt and keep the format of the corresponding text.

    Note that if the text is divided into two formats, it cannot be replaced

    Args:
        slide: slide object
        search_pattern: pattern to search
        repl: replacement string

    Returns:
        modified slide object
    """
    if repl is None:
        repl = ""
    search_pattern = re.compile(re.escape(search_str), re.IGNORECASE)
    for shape in slide.shapes:
        if shape.has_text_frame and not re.search(search_pattern, shape.text) is None:
            text_frame = shape.text_frame
            # shape.text_frame.text = re.sub(search_pattern, repl, shape.text_frame.text)
            # for m in re.finditer(search_pattern, text_frame.text):

            for paragraph in text_frame.paragraphs:
                # use style of the first run
                paragraph.runs[0].text = re.sub(search_pattern, repl, shape.text)
                # remove text in other runs
                for run in paragraph.runs[1:]:
                    run.text = ""
    return slide

def get_shape(
    slide: Slide,
    pic_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
    shape_type: Literal["picture", "table"] = "picture",
) -> BaseShape:
    if shape_type == "picture":
        _ele = [shape for shape in slide.shapes if isinstance(shape, Picture)]
    elif shape_type == "table":
        _ele = [shape for shape in slide.shapes if shape.has_table]

    if order == "t2b":
        ordered_shapes = sorted(
            _ele,
            key=lambda x: x.top,  # type: ignore
        )
    elif order == "l2r":
        ordered_shapes = sorted(
            _ele,
            key=lambda x: x.left,  # type: ignore
        )
    else:
        raise ValueError("order must be t2b or l2r")

    shape = ordered_shapes[pic_number]
    return shape


def replace_picture_in_slide(
    slide: Slide,
    fig: Union[
        bytes,
        str,
        plt.Figure,
        BinaryIO,
        alt.Chart,
        pd.DataFrame,
        Styler,
        List[List[Union[float, str]]],
    ],
    pic_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
    auto_reshape: bool = True,
) -> None:
    """替换某一页中PPT中的图片

    Args:
        slide: PPT页面对象
        fig: 用来替换的图片, 可以是从文件中读取的bytes, 可以是文件名, 可以是matplotlib的figure对象, 可以是文件的IO对象
        pic_number: 页面中第几个图片被替换, 按图片位置排序, 默认为0
        order: 替换图片的顺序, t2b表示从上到下, l2r表示从左到右, 默认为t2b
    """
    shape = get_shape(
        slide,
        pic_number,
        order,
        shape_type="picture",
    )
    replace_shape_with_picture(shape, fig, auto_reshape)


# def copy_row_insert_after(row:_RowCollection, copy_idx:int=-1, insert_idx:int=-1, init_cell_func: Optional[Callable[[_Cell], None]] = None):
#     '''
#     Duplicates target row to keep formatting and resets it's cells text_frames
#     (e.g. ``row = table.rows.copy_row_insert_after(0,1)``, which copies the first row and inserts after the second row as new third row).
#     Returns new |_Row| instance.
#     '''
#     new_row = copy.deepcopy(row._tbl.tr_lst[copy_idx])  # copies idx row element. Note: tr_lst idx is != _tbl idx.

#     for tc in new_row.tc_lst:
#         cell = _Cell(tc, new_row.tc_lst)
#         if init_cell_func:
#             init_cell_func(cell)

#     #_tbl[0] xml sets up the table and relationships <a:tblPr>: try table.rows.debug_tbl_idx(0)
#         #https://python-pptx.readthedocs.io/en/latest/dev/analysis/tbl-table.html?highlight=a%3AtblPr#xml-semantics
#     #_tbl[1] xml sets up the columns <a:tblGrid>: try table.rows.debug_tbl_idx(1)
#     #_tbl[2] xml is the first row <a:tr>: try table.rows.debug_tbl_idx(2)

#     self._tbl.insert(insert_idx, new_row) #rows begin starting idx 2. Need to read _tbl[0], _tbl[1] xml.

#     return _Row(new_row, self)


def replace_table_in_slide(
    slide: Slide,
    data: Union[
        pd.DataFrame,
        Styler,
        List[List[Union[float, str]]],
    ],
    shape_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
) -> None:
    shape = get_shape(
        slide,
        shape_number,
        order,
        shape_type="table",
    )
    if isinstance(data, List):
        df = pd.DataFrame(data)

    if isinstance(data, Styler):
        data._compute()  # type: ignore
        pandas_styles = data._translate(False, False)  # type: ignore
        df = data.data  # type: ignore
    x, y, cx, cy = (
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    # t = shape.table
    rn = len(df)
    cn = len(df[0])
    new_shape = slide.shapes.add_table(rn + 1, cn, x, y, cx, cy)

    # add headers
    for c in range(cn):
        new_shape.table.cell(0, c).text = html.unescape(
            pandas_styles["head"][0][c]["display_value"]
        )
    # add body
    for r in range(rn):
        for c in range(cn):
            # tc = copy.deepcopy(shape.table.cell(-1, -1)._tc)
            # new_shape.table.cell(r+1, c)._tc = tc
            new_shape.table.cell(r + 1, c).text = html.unescape(
                pandas_styles["body"][r][c]["display_value"]
            )
    old_shape = shape._element
    new_pic = new_shape._element
    old_shape.addnext(new_pic)
    old_shape.getparent().remove(old_shape)


def get_slide_from_shape(shape: BaseShape) -> Slide:
    for i in range(100):
        shape = shape._parent
        if isinstance(shape, Slide):
            return shape
    else:
        raise ValueError("shape is not in a slide")


def replace_shape_with_picture(
    shape: BaseShape,
    fig: Union[
        bytes,
        str,
        plt.Figure,
        BinaryIO,
        alt.Chart,
        pd.DataFrame,
        Styler,
        List[List[Union[float, str]]],
    ],
    auto_reshape: bool = True,
) -> None:
    # prepare figure
    if isinstance(fig, str) or isinstance(fig, IOBase):
        figio = fig
    elif isinstance(fig, bytes):
        figio = BytesIO(fig)
    elif isinstance(fig, plt.Figure):
        figio = BytesIO()
        if auto_reshape:
            fig.set_size_inches(shape.width.inches, shape.height.inches)
        fig.savefig(figio, format="png", bbox_inches="tight", dpi=160)
    elif isinstance(fig, alt.VegaLiteSchema):
        figio = BytesIO()
        if auto_reshape:
            fig = fig.properties(width=shape.width.pt, height=shape.height.pt)
        fig.save(figio, format="png", scale_factor=1.6)
    elif (
        isinstance(fig, List)
        or isinstance(fig, pd.DataFrame)
        or isinstance(fig, Styler)
    ):
        if isinstance(fig, List):
            fig = pd.DataFrame(fig)
        figio = BytesIO()
        dfi.export(fig, figio, dpi=160)
    else:
        raise ValueError(f"{type(fig)} {repr(fig)} is not supported")

    # replace picture
    slide = get_slide_from_shape(shape)
    new_shape = slide.shapes.add_picture(
        figio,
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    old_pic = shape._element
    new_pic = new_shape._element
    old_pic.addnext(new_pic)
    old_pic.getparent().remove(old_pic)

# ------------------------------------------------------
# PPT页面删除函数（对母版无影响）
# ------------------------------------------------------
def delete_slides_by_index(
    prs,        # pptx对象
    index_list  # 删除页码列表，为空则不删除
):
    slides = list(prs.slides._sldIdLst)
    for index in index_list:
        prs.slides._sldIdLst.remove(slides[index])

# ---------------------------------------------------------------------------
# PPT依照母版页码新增页面函数（对母版无影响）
# 新建的页面会自动添加到末尾，无法移动
# 利用此函数，母版只需保留互不相同的template页即可，通过复制增加的方式实现ppt的动态页数生成
# 返回：页面对象，可以直接进行图片插入等操作
# ---------------------------------------------------------------------------
def new_slide_by_index(
    prs,        # pptx对象
    index       # 添加新页面所依照的母版页码
):
    template = prs.slides[index]
    try:
        blank_slide_layout = prs.slide_layouts[index]
    except:
        blank_slide_layout = prs.slide_layouts[2]
    copied_slide = prs.slides.add_slide(blank_slide_layout)
    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel)
    return copied_slide

# ---------------------------------------------------------------------------
# 将dataframe styler对象拼接成一张图，输入styler构成的List
# ---------------------------------------------------------------------------
def combine_dataframe_styler_to_image(
    image_list  # Styler组成的list
):
    assert len(image_list)==2, '当前仅支持两个表拼接成一张图'
    df1 = image_list[0]
    df2 = image_list[1]
    assert isinstance(df1, Styler) and isinstance(df2, Styler), '列表中元素为Styler类型'
    f1 = BytesIO()
    f2 = BytesIO()
    result = BytesIO()
    dfi.export(df1, f1, dpi=160)
    dfi.export(df2, f2, dpi=160)
    image1 = Image.open(f1)
    image2 = Image.open(f2)
    total_height = image1.height + image2.height
    max_width = max(image1.width, image2.width)
    new_image = Image.new('RGB', (max_width, total_height))
    new_image.paste(image1, (0, 0))
    new_image.paste(image2, (0, image1.height))
    new_image.save(result, format='PNG')
    return result


# ---------------------------------------------------------------------------
# 将HTML的格式转化为图片，以BytesIO形式储存
# ---------------------------------------------------------------------------
def covert_html_to_image(
    html  # Styler组成的list
):

    filename = BytesIO()
    converter = dfi._pandas_accessor.Screenshot(
        fontsize=14,
        encode_base64=False,
        limit_crop=False,
        device_scale_factor=2
    ).run
    img_str = converter(html)
    filename.write(img_str)
    return filename
